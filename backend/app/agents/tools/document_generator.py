import logging
import json
import io
import base64
from typing import Any, Dict, List, Optional
from datetime import datetime

from app.agents.tools.agent_tool import AgentTool
from app.repositories.models.custom_bot import BotModel
from app.routes.schemas.conversation import type_model_name, DocumentToolResult
from pydantic import BaseModel, Field

# Document generation libraries
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from docx import Document
from docx.shared import Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)


class ExcelGeneratorInput(BaseModel):
    title: str = Field(description="Title for the Excel document")
    data: List[Dict[str, Any]] = Field(description="Data to populate the Excel sheet. Each dict represents a row with column names as keys.")
    sheet_name: str = Field(default="Sheet1", description="Name of the Excel sheet")
    include_header: bool = Field(default=True, description="Whether to include column headers")


class WordGeneratorInput(BaseModel):
    title: str = Field(description="Title for the Word document")
    content: List[Dict[str, Any]] = Field(description="Content sections. Each dict should have 'type' (heading, paragraph, list) and 'text' or 'items'")


class PowerPointGeneratorInput(BaseModel):
    title: str = Field(description="Title for the PowerPoint presentation")
    slides: List[Dict[str, Any]] = Field(description="Slide content. Each dict should have 'title' and 'content' (list of bullet points or paragraphs)")


def _generate_excel(tool_input: ExcelGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> DocumentToolResult:
    """Generate an Excel file from the provided data."""
    try:
        logger.info(f"Generating Excel document: {tool_input.title}")
        
        # Create workbook and worksheet
        wb = Workbook()
        ws = wb.active
        ws.title = tool_input.sheet_name
        
        # Add title
        ws['A1'] = tool_input.title
        ws['A1'].font = Font(bold=True, size=16)
        ws.merge_cells('A1:E1')
        ws['A1'].alignment = Alignment(horizontal='center')
        
        # Add data
        if tool_input.data:
            start_row = 3
            
            # Add headers if requested
            if tool_input.include_header and tool_input.data:
                headers = list(tool_input.data[0].keys())
                for col, header in enumerate(headers, 1):
                    cell = ws.cell(row=start_row, column=col, value=header)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(start_color="CCCCCC", end_color="CCCCCC", fill_type="solid")
                start_row += 1
            
            # Add data rows
            for row_idx, row_data in enumerate(tool_input.data, start_row):
                for col_idx, (key, value) in enumerate(row_data.items(), 1):
                    ws.cell(row=row_idx, column=col_idx, value=str(value))
        
        # Auto-adjust column widths
        for column in ws.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column_letter].width = adjusted_width
        
        # Save to bytes
        excel_buffer = io.BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        
        # Create filename
        filename = f"{tool_input.title.replace(' ', '_')}.xlsx"
        
        return DocumentToolResult(
            format="xlsx",
            name=filename,
            document=excel_buffer.getvalue()
        )
        
    except Exception as e:
        logger.error(f"Error generating Excel document: {e}")
        raise e


def _generate_word(tool_input: WordGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> DocumentToolResult:
    """Generate a Word document from the provided content."""
    try:
        logger.info(f"Generating Word document: {tool_input.title}")
        
        # Create document
        doc = Document()
        
        # Add title
        title_paragraph = doc.add_heading(tool_input.title, level=1)
        title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Add timestamp
        doc.add_paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph()  # Empty line
        
        # Add content sections
        for section in tool_input.content:
            section_type = section.get('type', 'paragraph')
            
            if section_type == 'heading':
                level = section.get('level', 2)
                doc.add_heading(section.get('text', ''), level=level)
                
            elif section_type == 'paragraph':
                doc.add_paragraph(section.get('text', ''))
                
            elif section_type == 'list':
                items = section.get('items', [])
                for item in items:
                    doc.add_paragraph(str(item), style='List Bullet')
                    
            elif section_type == 'table':
                table_data = section.get('data', [])
                if table_data:
                    table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                    table.style = 'Table Grid'
                    
                    for row_idx, row_data in enumerate(table_data):
                        for col_idx, cell_data in enumerate(row_data):
                            table.cell(row_idx, col_idx).text = str(cell_data)
        
        # Save to bytes
        word_buffer = io.BytesIO()
        doc.save(word_buffer)
        word_buffer.seek(0)
        
        # Create filename
        filename = f"{tool_input.title.replace(' ', '_')}.docx"
        
        return DocumentToolResult(
            format="docx",
            name=filename,
            document=word_buffer.getvalue()
        )
        
    except Exception as e:
        logger.error(f"Error generating Word document: {e}")
        raise e


def _generate_powerpoint(tool_input: PowerPointGeneratorInput, bot: BotModel | None, model: type_model_name | None) -> DocumentToolResult:
    """Generate a PowerPoint presentation from the provided slides."""
    try:
        logger.info(f"Generating PowerPoint presentation: {tool_input.title}")
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = tool_input.title
        title_slide.placeholders[1].text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Add content slides
        for slide_data in tool_input.slides:
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add slide title
            slide.shapes.title.text = slide_data.get('title', 'Slide')
            
            # Add content
            content = slide_data.get('content', [])
            if content:
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()
                
                for i, item in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    p.text = str(item)
                    p.level = 0
        
        # Save to bytes
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Create filename
        filename = f"{tool_input.title.replace(' ', '_')}.pptx"
        
        return DocumentToolResult(
            format="pptx",
            name=filename,
            document=ppt_buffer.getvalue()
        )
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint presentation: {e}")
        raise e


# Create the agent tools
excel_generator_tool = AgentTool(
    name="excel_generator",
    description="Generate an Excel spreadsheet with data. Useful for creating reports, data tables, and structured information.",
    args_schema=ExcelGeneratorInput,
    function=_generate_excel,
)

word_generator_tool = AgentTool(
    name="word_generator", 
    description="Generate a Word document with formatted content. Useful for creating reports, documentation, and text-based documents.",
    args_schema=WordGeneratorInput,
    function=_generate_word,
)

powerpoint_generator_tool = AgentTool(
    name="powerpoint_generator",
    description="Generate a PowerPoint presentation with slides. Useful for creating presentations, slide decks, and visual content.",
    args_schema=PowerPointGeneratorInput,
    function=_generate_powerpoint,
)